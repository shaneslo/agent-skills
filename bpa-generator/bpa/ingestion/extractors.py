"""1.2 Format Extractors — deterministic parse, dispatched on file extension.

PPT: python-pptx, no LLM. PDF: pdfplumber text + one cheap LLM classification
call to split ``pdf_procedure`` vs ``pdf_meeting`` (that label drives the
Normalizer's step-organization hint downstream).
"""
from __future__ import annotations

from pathlib import Path

from pydantic import BaseModel

from ..config import RunConfig
from ..models import get_model
from ..state import ExtractedDocument, GraphState, SourceDocumentType


# ---- 1.2a PPT extractor (deterministic) ----------------------------------- #
def _extract_ppt(path: str) -> ExtractedDocument:
    from pptx import Presentation

    prs = Presentation(path)
    slides: list[dict] = []
    text_parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        body: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            txt = shape.text_frame.text
            if shape == slide.shapes.title:
                title = txt
            elif txt.strip():
                body.append(txt)
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
        slides.append(
            {"slide_number": i, "title": title, "body": body, "speaker_notes": notes}
        )
        text_parts.append("\n".join([title, *body, notes]))
    return ExtractedDocument(
        source_document_type=SourceDocumentType.ppt,
        text="\n\n".join(text_parts),
        blocks=slides,
    )


# ---- 1.2b PDF extractor (deterministic text + 1 classification call) ------- #
class _PdfClass(BaseModel):
    label: str  # "pdf_procedure" | "pdf_meeting"


_PDF_CLASSIFY_SYSTEM = """\
You are classifying a PDF document as either:
- "pdf_procedure": structured procedural content with numbered headings, ordered
  steps, decision points, and minimal dialog
- "pdf_meeting": meeting notes or transcripts with attendee attribution,
  conversational flow, and actions described as discussion rather than steps
Examine the provided text. Output {"label": "pdf_procedure" | "pdf_meeting"}.
"""


def _extract_pdf(path: str, config: RunConfig) -> ExtractedDocument:
    import pdfplumber

    pages: list[dict] = []
    text_parts: list[str] = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            txt = page.extract_text() or ""
            pages.append({"page_number": i, "text": txt})
            text_parts.append(txt)
    full = "\n\n".join(text_parts)

    model = get_model("classify", config).with_structured_output(_PdfClass)
    result: _PdfClass = model.invoke(
        [("system", _PDF_CLASSIFY_SYSTEM), ("human", full[:6000])]
    )
    sub = (
        SourceDocumentType.pdf_meeting
        if result.label == "pdf_meeting"
        else SourceDocumentType.pdf_procedure
    )
    return ExtractedDocument(source_document_type=sub, text=full, blocks=pages)


def format_extractor(state: GraphState, config: RunConfig) -> GraphState:
    """Dispatch on extension; populate ``extracted`` and full ``raw_text``."""
    path = state["input_path"]
    suffix = Path(path).suffix.lower()
    if suffix in {".ppt", ".pptx"}:
        extracted = _extract_ppt(path)
    elif suffix == ".pdf":
        extracted = _extract_pdf(path, config)
    else:
        raise ValueError(f"Unsupported input format: {suffix!r} (expected .pptx/.pdf)")
    return {"extracted": extracted, "raw_text": extracted.text}
